#!/usr/bin/env python3
"""
MCP Windows Automation Server
==============================
Servidor MCP (Model Context Protocol) para controlar Windows automáticamente
desde Ollama o cualquier cliente MCP.

Autor: Freebuff AI
Licencia: MIT

Tools disponibles:
  - Ventanas:   list_windows, focus_window, move_window, resize_window, 
                minimize_window, maximize_window, close_window, get_active_window
  - Mouse:      mouse_move, mouse_click, mouse_scroll, mouse_position
  - Teclado:    keyboard_type, keyboard_hotkey, keyboard_press
  - Archivos:   file_list, file_read, file_write, file_delete, file_copy,
                folder_create, folder_list
  - Procesos:   process_list, process_kill, process_start
  - Sistema:    system_info, screenshot, clipboard_get, clipboard_set,
                open_url, run_command, volume_set, get_wifi_info
  - Registry:   registry_read
"""

import asyncio
import json
import os
import platform
import shutil
import socket
import subprocess
import sys
import time
import webbrowser
from datetime import datetime
from pathlib import Path

try:
    import pyautogui
    PYAUTOGUI_AVAILABLE = True
except ImportError:
    PYAUTOGUI_AVAILABLE = False

try:
    import pygetwindow as gw
    PYWINDOW_AVAILABLE = True
except ImportError:
    PYWINDOW_AVAILABLE = False

try:
    import psutil
    PSUTIL_AVAILABLE = True
except ImportError:
    PSUTIL_AVAILABLE = False

try:
    import pyperclip
    PYPERCLIP_AVAILABLE = True
except ImportError:
    PYPERCLIP_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    import win32gui
    import win32con
    import win32api
    import win32process
    import win32clipboard
    WIN32_AVAILABLE = True
except ImportError:
    WIN32_AVAILABLE = False

try:
    import winreg
    WINREG_AVAILABLE = True
except ImportError:
    WINREG_AVAILABLE = False

# ── Creación de documentos Office y PDF ─────────────────────────────────────
try:
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    from pptx.dml.color import RGBColor as PptxRGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from fpdf import FPDF
    FPDF_AVAILABLE = True
except ImportError:
    FPDF_AVAILABLE = False

try:
    import requests as _requests
    REQUESTS_AVAILABLE = True
except ImportError:
    REQUESTS_AVAILABLE = False

try:
    from markdownify import markdownify as _md
    MARKDOWNIFY_AVAILABLE = True
except ImportError:
    MARKDOWNIFY_AVAILABLE = False

# ─── MCP Imports ────────────────────────────────────────────────────────────
from mcp.server.fastmcp import FastMCP

# ─── Inicializar Servidor MCP ───────────────────────────────────────────────
mcp = FastMCP(
    "Windows Automation Server",
)

# ═════════════════════════════════════════════════════════════════════════════
# 0. GUARDIÁN ATLAS (modo guardián)
# Consulta atlas_guardian (E:\Agente_IA\memory_data\state\guardian.json) antes
# de operaciones sensibles. Si la config no existe, NO bloquea (modo abierto).
# Se puede apuntar a otra config con la env var ATLAS_GUARDIAN_CONFIG.
# ═════════════════════════════════════════════════════════════════════════════

_ATLAS_GUARDIAN_CANDIDATES = [
    os.environ.get("ATLAS_GUARDIAN_CONFIG", ""),
    r"E:\Agente_IA\memory_data\state\guardian.json",
    r"D:\Agente_IA\memory_data\state\guardian.json",
]

_guardian_cache = None  # (mtime, config)


def _guardian_config() -> dict:
    """Carga la config del guardián (cacheada por mtime)."""
    global _guardian_cache
    for cand in _ATLAS_GUARDIAN_CANDIDATES:
        if not cand:
            continue
        p = Path(cand)
        if p.exists():
            try:
                mtime = p.stat().st_mtime
                if _guardian_cache and _guardian_cache[0] == mtime:
                    return _guardian_cache[1]
                cfg = json.loads(p.read_text(encoding="utf-8"))
                _guardian_cache = (mtime, cfg)
                return cfg
            except Exception:
                return {}
    return {}


def _guardian_decision(operation: str, params: dict) -> dict:
    """
    Devuelve {allowed: bool, reason: str}.
    - Sin config → allowed (no bloquea).
    - relax → allowed.
    - strict → bloquea run_script/process_kill/registry_write.
    - lista blanca de binarios/procesos/rutas según operación.
    """
    cfg = _guardian_config()
    if not cfg:
        return {"allowed": True, "reason": "sin config guardián (modo abierto)"}
    level = str(cfg.get("level", "guard")).lower()

    if level == "relax":
        return {"allowed": True, "reason": "modo relax"}

    if level == "strict" and operation in cfg.get("blocked_ops", []):
        return {"allowed": False, "reason": f"operación '{operation}' bloqueada en modo strict"}

    wl_bin = [b.lower() for b in cfg.get("whitelist_binaries", [])]
    wl_proc = [p.lower() for p in cfg.get("whitelist_processes", [])]

    if operation in ("run_command", "process_start"):
        cmd = params.get("command", "") or ""
        binary = cmd.split()[0].strip().lower()
        if binary and not any(b in binary for b in wl_bin):
            return {"allowed": False, "reason": f"binario '{binary}' no está en lista blanca"}

    elif operation == "run_script":
        sp = params.get("script_path", "") or ""
        binary = Path(sp).stem.lower() if sp else ""
        if sp and binary and not any(b in binary for b in wl_bin):
            return {"allowed": False, "reason": f"script '{binary}' no está en lista blanca"}

    elif operation == "process_kill":
        if level == "guard":
            return {"allowed": True, "reason": "confirmación requerida (modo guard)"}

    elif operation == "file_delete":
        path = params.get("path", "") or ""
        if path:
            try:
                p = Path(path).resolve()
                allowed_dirs = [str(Path(a).resolve()) for a in cfg.get("allowed_dirs", [])]
                if not any(p.is_relative_to(a) for a in allowed_dirs):
                    return {"allowed": False, "reason": f"ruta '{path}' fuera de directorios permitidos"}
            except Exception:
                return {"allowed": False, "reason": f"no se pudo validar ruta '{path}'"}

    return {"allowed": True, "reason": "ok"}


def _guardian_gate(operation: str, params: dict) -> str | None:
    """
    Puerta del guardián. Devuelve el JSON de bloqueo (string) si la operación
    NO está permitida, o None si puede ejecutarse.
    """
    d = _guardian_decision(operation, params)
    if not d["allowed"]:
        return json.dumps({
            "error": "BLOQUEADO por atlas-guardian",
            "guard_block": True,
            "operation": operation,
            "reason": d["reason"],
        }, ensure_ascii=False, indent=2)
    return None


# ═════════════════════════════════════════════════════════════════════════════
# 1. GESTIÓN DE VENTANAS
# ═════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def list_windows() -> str:
    """
    Lista todas las ventanas abiertas visibles en Windows.
    
    Returns:
        JSON string con lista de ventanas (título, handle, posición, tamaño)
    """
    if not PYWINDOW_AVAILABLE and not WIN32_AVAILABLE:
        return json.dumps({"error": "Se requiere pygetwindow o pywin32"})
    
    windows = []
    
    if PYWINDOW_AVAILABLE:
        all_windows = gw.getAllWindows()
        for win in all_windows:
            if win.title.strip():
                windows.append({
                    "title": win.title,
                    "left": win.left,
                    "top": win.top,
                    "width": win.width,
                    "height": win.height,
                    "is_active": win.isActive,
                    "is_minimized": win.isMinimized,
                    "is_maximized": win.isMaximized,
                })
    elif WIN32_AVAILABLE:
        def enum_callback(hwnd, _):
            if win32gui.IsWindowVisible(hwnd):
                title = win32gui.GetWindowText(hwnd)
                if title:
                    rect = win32gui.GetWindowRect(hwnd)
                    windows.append({
                        "title": title,
                        "hwnd": hwnd,
                        "left": rect[0],
                        "top": rect[1],
                        "right": rect[2],
                        "bottom": rect[3],
                    })
        win32gui.EnumWindows(enum_callback, None)
    
    return json.dumps(windows, ensure_ascii=False, indent=2)


@mcp.tool()
def focus_window(title: str) -> str:
    """
    Enfoca (activa) una ventana por su título.
    
    Args:
        title: Título (o parte del título) de la ventana a enfocar
    """
    if not PYWINDOW_AVAILABLE:
        return json.dumps({"error": "Se requiere pygetwindow"})
    
    try:
        windows = gw.getWindowsWithTitle(title)
        if not windows:
            return json.dumps({"error": f"No se encontró ventana con título: {title}"})
        
        win = windows[0]
        win.activate()
        time.sleep(0.3)
        return json.dumps({"success": True, "title": win.title, "action": "focus"})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def move_window(title: str, x: int, y: int) -> str:
    """
    Mueve una ventana a una posición específica en la pantalla.
    
    Args:
        title: Título (o parte) de la ventana
        x: Posición horizontal en píxeles
        y: Posición vertical en píxeles
    """
    if not PYWINDOW_AVAILABLE:
        return json.dumps({"error": "Se requiere pygetwindow"})
    
    try:
        windows = gw.getWindowsWithTitle(title)
        if not windows:
            return json.dumps({"error": f"No se encontró ventana con título: {title}"})
        
        win = windows[0]
        win.moveTo(x, y)
        return json.dumps({"success": True, "title": win.title, "x": x, "y": y})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def resize_window(title: str, width: int, height: int) -> str:
    """
    Redimensiona una ventana.
    
    Args:
        title: Título (o parte) de la ventana
        width: Nuevo ancho en píxeles
        height: Nuevo alto en píxeles
    """
    if not PYWINDOW_AVAILABLE:
        return json.dumps({"error": "Se requiere pygetwindow"})
    
    try:
        windows = gw.getWindowsWithTitle(title)
        if not windows:
            return json.dumps({"error": f"No se encontró ventana con título: {title}"})
        
        win = windows[0]
        win.resizeTo(width, height)
        return json.dumps({"success": True, "title": win.title, "width": width, "height": height})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def minimize_window(title: str) -> str:
    """Minimiza una ventana."""
    if not PYWINDOW_AVAILABLE:
        return json.dumps({"error": "Se requiere pygetwindow"})
    try:
        windows = gw.getWindowsWithTitle(title)
        if not windows:
            return json.dumps({"error": f"No se encontró ventana: {title}"})
        windows[0].minimize()
        return json.dumps({"success": True, "action": "minimize", "title": windows[0].title})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def maximize_window(title: str) -> str:
    """Maximiza una ventana."""
    if not PYWINDOW_AVAILABLE:
        return json.dumps({"error": "Se requiere pygetwindow"})
    try:
        windows = gw.getWindowsWithTitle(title)
        if not windows:
            return json.dumps({"error": f"No se encontró ventana: {title}"})
        windows[0].maximize()
        return json.dumps({"success": True, "action": "maximize", "title": windows[0].title})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def close_window(title: str) -> str:
    """Cierra una ventana."""
    if not PYWINDOW_AVAILABLE:
        return json.dumps({"error": "Se requiere pygetwindow"})
    try:
        windows = gw.getWindowsWithTitle(title)
        if not windows:
            return json.dumps({"error": f"No se encontró ventana: {title}"})
        windows[0].close()
        return json.dumps({"success": True, "action": "close", "title": title})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def get_active_window() -> str:
    """Obtiene información de la ventana activa actual."""
    result = {}
    
    if WIN32_AVAILABLE:
        hwnd = win32gui.GetForegroundWindow()
        title = win32gui.GetWindowText(hwnd)
        rect = win32gui.GetWindowRect(hwnd)
        result = {
            "title": title,
            "hwnd": hwnd,
            "left": rect[0],
            "top": rect[1],
            "right": rect[2],
            "bottom": rect[3],
        }
    elif PYWINDOW_AVAILABLE:
        try:
            win = gw.getActiveWindow()
            if win:
                result = {
                    "title": win.title,
                    "left": win.left,
                    "top": win.top,
                    "width": win.width,
                    "height": win.height,
                }
        except Exception:
            result = {"error": "No se pudo obtener la ventana activa"}
    
    return json.dumps(result, ensure_ascii=False, indent=2)


# ═════════════════════════════════════════════════════════════════════════════
# 2. CONTROL DE MOUSE
# ═════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def mouse_move(x: int, y: int, duration: float = 0.5) -> str:
    """
    Mueve el cursor del mouse a una posición específica.
    
    Args:
        x: Coordenada X
        y: Coordenada Y
        duration: Duración del movimiento en segundos (default: 0.5)
    """
    if not PYAUTOGUI_AVAILABLE:
        return json.dumps({"error": "Se requiere pyautogui"})
    
    try:
        pyautogui.moveTo(x, y, duration=duration)
        return json.dumps({"success": True, "x": x, "y": y})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def mouse_click(x: int = None, y: int = None, button: str = "left", 
                clicks: int = 1) -> str:
    """
    Hace click en una posición específica o en la posición actual.
    
    Args:
        x: Coordenada X (opcional, usa posición actual si no se especifica)
        y: Coordenada Y (opcional)
        button: Botón del mouse ('left', 'right', 'middle')
        clicks: Número de clicks (1=simple, 2=doble)
    """
    if not PYAUTOGUI_AVAILABLE:
        return json.dumps({"error": "Se requiere pyautogui"})
    
    try:
        if x is not None and y is not None:
            pyautogui.click(x, y, clicks=clicks, button=button)
            return json.dumps({"success": True, "x": x, "y": y, "button": button, "clicks": clicks})
        else:
            pyautogui.click(clicks=clicks, button=button)
            return json.dumps({"success": True, "button": button, "clicks": clicks})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def mouse_scroll(clicks: int, x: int = None, y: int = None) -> str:
    """
    Hace scroll con el mouse.
    
    Args:
        clicks: Número de clicks de scroll (positivo=arriba, negativo=abajo)
        x: Posición X opcional
        y: Posición Y opcional
    """
    if not PYAUTOGUI_AVAILABLE:
        return json.dumps({"error": "Se requiere pyautogui"})
    
    try:
        if x is not None and y is not None:
            pyautogui.scroll(clicks, x, y)
        else:
            pyautogui.scroll(clicks)
        return json.dumps({"success": True, "scroll_clicks": clicks})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def mouse_position() -> str:
    """Obtiene la posición actual del cursor del mouse."""
    if not PYAUTOGUI_AVAILABLE:
        return json.dumps({"error": "Se requiere pyautogui"})
    
    try:
        x, y = pyautogui.position()
        return json.dumps({"x": x, "y": y})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def mouse_drag(x: int, y: int, button: str = "left", duration: float = 0.5) -> str:
    """
    Arrastra el mouse (click sostenido + movimiento).
    
    Args:
        x: Coordenada X destino
        y: Coordenada Y destino
        button: Botón ('left', 'right', 'middle')
        duration: Duración del arrastre en segundos
    """
    if not PYAUTOGUI_AVAILABLE:
        return json.dumps({"error": "Se requiere pyautogui"})
    
    try:
        pyautogui.dragTo(x, y, duration=duration, button=button)
        return json.dumps({"success": True, "x": x, "y": y, "button": button})
    except Exception as e:
        return json.dumps({"error": str(e)})


# ═════════════════════════════════════════════════════════════════════════════
# 3. CONTROL DE TECLADO
# ═════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def keyboard_type(text: str, interval: float = 0.05) -> str:
    """
    Escribe texto en la ventana activa.
    
    Args:
        text: Texto a escribir
        interval: Intervalo entre caracteres en segundos (default: 0.05)
    """
    if not PYAUTOGUI_AVAILABLE:
        return json.dumps({"error": "Se requiere pyautogui"})
    
    try:
        pyautogui.typewrite(text, interval=interval)
        return json.dumps({"success": True, "chars": len(text)})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def keyboard_hotkey(keys: str) -> str:
    """
    Ejecuta un atajo de teclado (ej: 'ctrl+c' para copiar, 'alt+tab' para cambiar ventana).
    
    Args:
        keys: Combinación de teclas separadas por '+' (ej: 'ctrl+c', 'alt+tab', 'win+d')
    """
    if not PYAUTOGUI_AVAILABLE:
        return json.dumps({"error": "Se requiere pyautogui"})
    
    try:
        key_list = [k.strip() for k in keys.split("+")]
        pyautogui.hotkey(*key_list)
        return json.dumps({"success": True, "keys": key_list})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def keyboard_press(key: str, presses: int = 1) -> str:
    """
    Presiona una tecla específica.
    
    Args:
        key: Nombre de la tecla (ej: 'enter', 'tab', 'escape', 'f5')
        presses: Número de veces a presionar
    """
    if not PYAUTOGUI_AVAILABLE:
        return json.dumps({"error": "Se requiere pyautogui"})
    
    try:
        pyautogui.press(key, presses=presses)
        return json.dumps({"success": True, "key": key, "presses": presses})
    except Exception as e:
        return json.dumps({"error": str(e)})


# ═════════════════════════════════════════════════════════════════════════════
# 4. OPERACIONES DE ARCHIVOS
# ═════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def file_list(path: str = ".") -> str:
    """
    Lista archivos y directorios en una ruta.
    
    Args:
        path: Ruta del directorio (default: directorio actual)
    """
    try:
        p = Path(path).expanduser().resolve()
        if not p.exists():
            return json.dumps({"error": f"La ruta no existe: {path}"})
        if not p.is_dir():
            return json.dumps({"error": f"No es un directorio: {path}"})
        
        items = []
        for item in p.iterdir():
            items.append({
                "name": item.name,
                "type": "directory" if item.is_dir() else "file",
                "size": item.stat().st_size if not item.is_dir() else 0,
                "modified": datetime.fromtimestamp(item.stat().st_mtime).isoformat(),
            })
        
        return json.dumps({
            "path": str(p),
            "count": len(items),
            "items": items,
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def file_read(path: str) -> str:
    """
    Lee el contenido de un archivo de texto.
    
    Args:
        path: Ruta del archivo
    """
    try:
        p = Path(path).expanduser().resolve()
        if not p.exists():
            return json.dumps({"error": f"El archivo no existe: {path}"})
        if not p.is_file():
            return json.dumps({"error": f"No es un archivo: {path}"})
        
        content = p.read_text(encoding="utf-8", errors="replace")
        return json.dumps({
            "path": str(p),
            "size": len(content),
            "content": content,
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def file_write(path: str, content: str) -> str:
    """
    Escribe contenido en un archivo (lo crea si no existe).
    
    Args:
        path: Ruta del archivo
        content: Contenido a escribir
    """
    try:
        p = Path(path).expanduser().resolve()
        existed_before = p.exists()
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_text(content, encoding="utf-8")
        return json.dumps({
            "success": True,
            "path": str(p),
            "size": len(content),
            "action": "updated" if existed_before else "created",
        }, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def file_delete(path: str) -> str:
    """
    Elimina un archivo o directorio.
    
    Args:
        path: Ruta del archivo o directorio a eliminar
    """
    block = _guardian_gate("file_delete", {"path": path})
    if block:
        return block
    try:
        p = Path(path).expanduser().resolve()
        if not p.exists():
            return json.dumps({"error": f"No existe: {path}"})
        
        if p.is_dir():
            shutil.rmtree(p)
            return json.dumps({"success": True, "action": "deleted_directory", "path": str(p)})
        else:
            p.unlink()
            return json.dumps({"success": True, "action": "deleted_file", "path": str(p)})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def file_copy(source: str, destination: str) -> str:
    """
    Copia un archivo o directorio.
    
    Args:
        source: Ruta origen
        destination: Ruta destino
    """
    try:
        src = Path(source).expanduser().resolve()
        dst = Path(destination).expanduser().resolve()
        
        if not src.exists():
            return json.dumps({"error": f"El origen no existe: {source}"})
        
        if src.is_dir():
            shutil.copytree(src, dst)
        else:
            dst.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(src, dst)
        
        return json.dumps({"success": True, "from": str(src), "to": str(dst)})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def folder_create(path: str) -> str:
    """
    Crea un directorio (y sus padres si es necesario).
    
    Args:
        path: Ruta del directorio a crear
    """
    try:
        p = Path(path).expanduser().resolve()
        p.mkdir(parents=True, exist_ok=True)
        return json.dumps({"success": True, "path": str(p)})
    except Exception as e:
        return json.dumps({"error": str(e)})


# ═════════════════════════════════════════════════════════════════════════════
# 5. GESTIÓN DE PROCESOS
# ═════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def process_list() -> str:
    """
    Lista los procesos en ejecución.
    
    Returns:
        JSON con lista de procesos (PID, nombre, CPU%, memoria MB)
    """
    if not PSUTIL_AVAILABLE:
        return json.dumps({"error": "Se requiere psutil"})
    
    try:
        processes = []
        for proc in psutil.process_iter(['pid', 'name', 'cpu_percent', 'memory_info', 'status']):
            try:
                info = proc.info
                processes.append({
                    "pid": info['pid'],
                    "name": info['name'],
                    "cpu_percent": info['cpu_percent'],
                    "memory_mb": round(info['memory_info'].rss / (1024 * 1024), 2) if info['memory_info'] else 0,
                    "status": info['status'],
                })
            except (psutil.NoSuchProcess, psutil.AccessDenied):
                pass
        
        # Ordenar por uso de memoria descendente
        processes.sort(key=lambda p: p['memory_mb'], reverse=True)
        
        return json.dumps({
            "count": len(processes),
            "processes": processes[:100],  # Limitar a 100 procesos
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def process_kill(pid: int) -> str:
    """
    Mata un proceso por su PID.
    
    Args:
        pid: PID del proceso a matar
    """
    block = _guardian_gate("process_kill", {"pid": pid})
    if block:
        return block
    if not PSUTIL_AVAILABLE:
        return json.dumps({"error": "Se requiere psutil"})
    
    try:
        proc = psutil.Process(pid)
        name = proc.name()
        proc.terminate()
        return json.dumps({"success": True, "pid": pid, "name": name, "action": "terminated"})
    except psutil.NoSuchProcess:
        return json.dumps({"error": f"No existe proceso con PID {pid}"})
    except psutil.AccessDenied:
        return json.dumps({"error": f"No tienes permisos para matar el proceso {pid}"})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def process_start(command: str, args: str = "") -> str:
    """
    Inicia un proceso/programa.
    
    Args:
        command: Comando o ruta del ejecutable
        args: Argumentos adicionales (opcional)
    """
    block = _guardian_gate("process_start", {"command": command})
    if block:
        return block
    try:
        full_cmd = f"{command} {args}".strip()
        subprocess.Popen(full_cmd, shell=True)
        return json.dumps({"success": True, "command": full_cmd})
    except Exception as e:
        return json.dumps({"error": str(e)})


# ═════════════════════════════════════════════════════════════════════════════
# 6. INFORMACIÓN DEL SISTEMA
# ═════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def system_info() -> str:
    """Obtiene información detallada del sistema Windows."""
    info = {
        "system": platform.system(),
        "version": platform.version(),
        "release": platform.release(),
        "machine": platform.machine(),
        "processor": platform.processor(),
        "hostname": socket.gethostname(),
        "python_version": sys.version,
    }
    
    if PSUTIL_AVAILABLE:
        import psutil
        mem = psutil.virtual_memory()
        disk = psutil.disk_usage('/')
        info["memory"] = {
            "total_gb": round(mem.total / (1024**3), 2),
            "available_gb": round(mem.available / (1024**3), 2),
            "percent_used": mem.percent,
        }
        info["disk"] = {
            "total_gb": round(disk.total / (1024**3), 2),
            "free_gb": round(disk.free / (1024**3), 2),
            "percent_used": disk.percent,
        }
        info["cpu"] = {
            "cores": psutil.cpu_count(logical=False),
            "logical_cores": psutil.cpu_count(logical=True),
            "percent": psutil.cpu_percent(interval=0.5),
        }
        info["boot_time"] = datetime.fromtimestamp(psutil.boot_time()).isoformat()
    
    return json.dumps(info, ensure_ascii=False, indent=2)


@mcp.tool()
def screenshot() -> str:
    """
    Toma una captura de pantalla y la guarda como archivo PNG.
    
    Returns:
        Ruta del archivo de captura
    """
    if not PYAUTOGUI_AVAILABLE:
        return json.dumps({"error": "Se requiere pyautogui"})
    
    try:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"screenshot_{timestamp}.png"
        filepath = str(Path.home() / "Pictures" / filename)
        
        # Asegurar que el directorio existe
        Path(filepath).parent.mkdir(parents=True, exist_ok=True)
        
        screenshot_img = pyautogui.screenshot()
        screenshot_img.save(filepath)
        
        return json.dumps({
            "success": True,
            "filepath": filepath,
            "size": f"{screenshot_img.width}x{screenshot_img.height}",
        }, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def clipboard_get() -> str:
    """Obtiene el contenido actual del portapapeles."""
    if PYPERCLIP_AVAILABLE:
        try:
            content = pyperclip.paste()
            return json.dumps({"content": content, "length": len(content)}, ensure_ascii=False)
        except Exception:
            pass
    
    if WIN32_AVAILABLE:
        try:
            win32clipboard.OpenClipboard()
            data = win32clipboard.GetClipboardData()
            win32clipboard.CloseClipboard()
            return json.dumps({"content": data, "length": len(data)}, ensure_ascii=False)
        except Exception as e:
            return json.dumps({"error": str(e)})
    
    return json.dumps({"error": "No hay módulo para portapapeles disponible"})


@mcp.tool()
def clipboard_set(content: str) -> str:
    """
    Establece el contenido del portapapeles.
    
    Args:
        content: Texto a copiar al portapapeles
    """
    if PYPERCLIP_AVAILABLE:
        try:
            pyperclip.copy(content)
            return json.dumps({"success": True, "length": len(content)})
        except Exception:
            pass
    
    if WIN32_AVAILABLE:
        try:
            win32clipboard.OpenClipboard()
            win32clipboard.EmptyClipboard()
            win32clipboard.SetClipboardText(content)
            win32clipboard.CloseClipboard()
            return json.dumps({"success": True, "length": len(content)})
        except Exception as e:
            return json.dumps({"error": str(e)})
    
    return json.dumps({"error": "No hay módulo para portapapeles disponible"})


@mcp.tool()
def open_url(url: str) -> str:
    """
    Abre una URL en el navegador predeterminado.
    
    Args:
        url: URL a abrir (ej: https://www.google.com)
    """
    try:
        webbrowser.open(url)
        return json.dumps({"success": True, "url": url})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def run_command(command: str, timeout: int = 30) -> str:
    """
    Ejecuta un comando en CMD y devuelve la salida.
    
    Args:
        command: Comando a ejecutar
        timeout: Timeout en segundos (default: 30)
    """
    block = _guardian_gate("run_command", {"command": command})
    if block:
        return block
    try:
        result = subprocess.run(
            command,
            shell=True,
            capture_output=True,
            text=True,
            timeout=timeout,
        )
        return json.dumps({
            "success": result.returncode == 0,
            "returncode": result.returncode,
            "stdout": result.stdout,
            "stderr": result.stderr,
        }, ensure_ascii=False, indent=2)
    except subprocess.TimeoutExpired:
        return json.dumps({"error": f"Comando agotó el tiempo de espera ({timeout}s)"})
    except Exception as e:
        return json.dumps({"error": str(e)})

@mcp.tool()
def run_script(script_path: str, args: str = "") -> str:
    """
    Ejecuta un script externo (Batch .bat o PowerShell .ps1).
    Útil para workflows de automatización complejos.
    
    Args:
        script_path: Ruta al archivo de script
        args: Argumentos para el script
    """
    block = _guardian_gate("run_script", {"script_path": script_path})
    if block:
        return block
    try:
        p = Path(script_path).expanduser().resolve()
        if not p.exists():
            return json.dumps({"error": f"El script no existe: {script_path}"})
        
        ext = p.suffix.lower()
        if ext == ".ps1":
            cmd = ["powershell", "-ExecutionPolicy", "Bypass", "-File", str(p)]
            if args:
                cmd.extend(args.split())
        elif ext == ".bat" or ext == ".cmd":
            cmd = [str(p)]
            if args:
                cmd.extend(args.split())
        else:
            return json.dumps({"error": "Solo se admiten scripts .ps1, .bat o .cmd"})

        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            shell=(ext in (".bat", ".cmd")),
            timeout=60
        )
        return json.dumps({
            "success": result.returncode == 0,
            "returncode": result.returncode,
            "stdout": result.stdout,
            "stderr": result.stderr,
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def volume_set(level: int) -> str:
    """
    Establece el volumen del sistema (0-100).
    
    Args:
        level: Nivel de volumen (0 a 100)
    """
    if not (0 <= level <= 100):
        return json.dumps({"error": "El nivel debe estar entre 0 y 100"})
    
    try:
        # Usar PowerShell para ajustar volumen (más compatible)
        cmd = f'(New-Object -ComObject WScript.Shell).SendKeys([char]174); ' \
              f'for($i=0;$i -lt 50;$i++){{' \
              f'  (New-Object -ComObject WScript.Shell).SendKeys([char]175)' \
              f'}}'
        
        # Alternativa usando native Windows API via PowerShell
        ps_cmd = f'''
        $obj = New-Object -ComObject WScript.Shell
        for($i=0; $i -lt 100; $i++) {{ $obj.SendKeys([char]174) }}
        for($i=0; $i -lt {level}; $i++) {{ $obj.SendKeys([char]175) }}
        '''
        subprocess.run(["powershell", "-Command", ps_cmd], capture_output=True, timeout=5)
        
        return json.dumps({"success": True, "volume": level})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def get_wifi_info() -> str:
    """Obtiene información de la red WiFi actual."""
    try:
        # Obtener perfil WiFi actual
        result = subprocess.run(
            ["netsh", "wlan", "show", "interfaces"],
            capture_output=True, text=True, timeout=10
        )
        
        # Obtener todas las redes guardadas
        profiles_result = subprocess.run(
            ["netsh", "wlan", "show", "profiles"],
            capture_output=True, text=True, timeout=10
        )
        
        return json.dumps({
            "current_interface": result.stdout,
            "saved_profiles": profiles_result.stdout,
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)})


# ═════════════════════════════════════════════════════════════════════════════
# 7. REGISTRO DE WINDOWS
# ═════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def registry_read(key_path: str, value_name: str = None) -> str:
    """
    Lee un valor del registro de Windows.
    
    Args:
        key_path: Ruta de la clave en el registro (ej: SOFTWARE\\Microsoft\\Windows\\CurrentVersion)
        value_name: Nombre del valor a leer (opcional, si es None lista los valores)
    """
    if not WINREG_AVAILABLE:
        return json.dumps({"error": "winreg no está disponible"})
    
    try:
        hive_map = {
            "HKEY_LOCAL_MACHINE": winreg.HKEY_LOCAL_MACHINE,
            "HKLM": winreg.HKEY_LOCAL_MACHINE,
            "HKEY_CURRENT_USER": winreg.HKEY_CURRENT_USER,
            "HKCU": winreg.HKEY_CURRENT_USER,
            "HKEY_CLASSES_ROOT": winreg.HKEY_CLASSES_ROOT,
            "HKCR": winreg.HKEY_CLASSES_ROOT,
            "HKEY_USERS": winreg.HKEY_USERS,
            "HKU": winreg.HKEY_USERS,
        }
        
        # Parsear hive
        parts = key_path.split("\\", 1)
        hive_str = parts[0].upper()
        subkey = parts[1] if len(parts) > 1 else ""
        
        if hive_str in hive_map:
            hive = hive_map[hive_str]
        else:
            hive = winreg.HKEY_CURRENT_USER
            subkey = key_path
        
        with winreg.OpenKey(hive, subkey) as key:
            if value_name:
                value, reg_type = winreg.QueryValueEx(key, value_name)
                return json.dumps({
                    "key": key_path,
                    "value_name": value_name,
                    "value": str(value),
                    "type": reg_type,
                }, ensure_ascii=False)
            else:
                values = []
                i = 0
                while True:
                    try:
                        name, value, reg_type = winreg.EnumValue(key, i)
                        values.append({
                            "name": name,
                            "value": str(value)[:200],  # Limitar longitud
                            "type": reg_type,
                        })
                        i += 1
                    except OSError:
                        break
                
                subkeys = []
                i = 0
                while True:
                    try:
                        subkeys.append(winreg.EnumKey(key, i))
                        i += 1
                    except OSError:
                        break
                
                return json.dumps({
                    "key": key_path,
                    "values": values,
                    "subkeys": subkeys,
                }, ensure_ascii=False, indent=2)
    except FileNotFoundError:
        return json.dumps({"error": f"Clave no encontrada: {key_path}"})
    except Exception as e:
        return json.dumps({"error": str(e)})


# ═════════════════════════════════════════════════════════════════════════════
# 8. NOTIFICACIONES
# ═════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def notify(title: str, message: str, duration: int = 5) -> str:
    """
    Muestra una notificación del sistema en Windows.
    
    Args:
        title: Título de la notificación
        message: Mensaje de la notificación
        duration: Duración en segundos (default: 5)
    """
    try:
        # Usar PowerShell para mostrar notificación tipo Toast
        ps_script = f'''
        [Windows.UI.Notifications.ToastNotificationManager, Windows.UI.Notifications, ContentType = WindowsRuntime] > $null
        $template = [Windows.UI.Notifications.ToastNotificationManager]::GetTemplateContent([Windows.UI.Notifications.ToastTemplateType]::ToastText02)
        $textNodes = $template.GetElementsByTagName("text")
        $textNodes.Item(0).AppendChild($template.CreateTextNode("{title}")) > $null
        $textNodes.Item(1).AppendChild($template.CreateTextNode("{message}")) > $null
        $toast = [Windows.UI.Notifications.ToastNotification]::new($template)
        [Windows.UI.Notifications.ToastNotificationManager]::CreateToastNotifier("MCP Windows AI").Show($toast)
        '''
        
        subprocess.run(["powershell", "-Command", ps_script], capture_output=True, timeout=10)
        return json.dumps({"success": True, "title": title, "message": message})
    except Exception as e:
        # Fallback: usar msg.exe
        try:
            subprocess.run(["msg", "*", f"{title}: {message}"], capture_output=True, timeout=5)
            return json.dumps({"success": True, "fallback": "msg.exe"})
        except Exception:
            return json.dumps({"error": str(e)})


# ═════════════════════════════════════════════════════════════════════════════
# 9. CREACIÓN DE DOCUMENTOS (WORD, EXCEL, POWERPOINT, PDF)
# ═════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def docx_create(path: str, content: str) -> str:
    """
    Crea un documento Word (.docx) con formato profesional.
    
    Args:
        path: Ruta del archivo .docx a crear
        content: JSON string con la estructura del documento:
            {
              "title": "Título principal",
              "subtitle": "Subtítulo (opcional)",
              "author": "Autor (opcional)",
              "sections": [
                {"heading": "Sección 1", "level": 1,
                 "paragraphs": ["Texto del párrafo 1", "Párrafo 2"],
                 "bullets": ["punto 1", "punto 2"],
                 "table": {"headers": ["Col1","Col2"], "rows": [["a","b"],["c","d"]]}
                }
              ]
            }
    """
    if not DOCX_AVAILABLE:
        return json.dumps({"error": "pip install python-docx"})
    
    try:
        data = json.loads(content) if isinstance(content, str) else content
        p = Path(path).expanduser().resolve()
        p.parent.mkdir(parents=True, exist_ok=True)
        
        doc = Document()
        
        # Título
        if data.get("title"):
            doc.add_heading(data["title"], level=0)
        if data.get("subtitle"):
            sub = doc.add_paragraph(data["subtitle"])
            sub.runs[0].italic = True if sub.runs else None
            sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if data.get("author"):
            auth = doc.add_paragraph(f"Por: {data['author']}")
            auth.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Secciones
        for section in data.get("sections", []):
            if section.get("heading"):
                doc.add_heading(section["heading"], level=section.get("level", 1))
            
            for para in section.get("paragraphs", []):
                doc.add_paragraph(para)
            
            for bullet in section.get("bullets", []):
                doc.add_paragraph(bullet, style="List Bullet")
            
            for num in section.get("numbered", []):
                doc.add_paragraph(num, style="List Number")
            
            table_data = section.get("table")
            if table_data:
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                if headers:
                    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
                    table.style = "Light Grid Accent 1"
                    for i, h in enumerate(headers):
                        cell = table.rows[0].cells[i]
                        cell.text = str(h)
                        for run in cell.paragraphs[0].runs:
                            run.bold = True
                    for r, row in enumerate(rows):
                        for c, val in enumerate(row):
                            if c < len(headers):
                                table.rows[r + 1].cells[c].text = str(val)
        
        doc.save(str(p))
        return json.dumps({"success": True, "path": str(p), "type": "docx"}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": f"Error creando Word: {e}"})


@mcp.tool()
def xlsx_create(path: str, content: str) -> str:
    """
    Crea un archivo Excel (.xlsx) con datos, FÓRMULAS y formato.
    
    Args:
        path: Ruta del archivo .xlsx a crear
        content: JSON string con la estructura:
            {
              "sheets": [
                {
                  "name": "Ventas",
                  "headers": ["Producto", "Precio", "Cantidad", "Total"],
                  "rows": [
                    ["Laptop", 1000, 5, "=B2*C2"],
                    ["Mouse", 25, 10, "=B3*C3"]
                  ],
                  "formulas": {"D10": "=SUM(D2:D3)", "B12": "=AVERAGE(B2:B3)"},
                  "column_widths": {"A": 20, "B": 12},
                  "header_color": "4472C4"
                }
              ]
            }
            Las fórmulas usan sintaxis de Excel: =SUM(), =AVERAGE(), =VLOOKUP(), =SI(), etc.
            Las celdas en "rows" que empiecen con "=" se guardan como fórmula.
    """
    if not XLSX_AVAILABLE:
        return json.dumps({"error": "pip install openpyxl"})
    
    try:
        data = json.loads(content) if isinstance(content, str) else content
        p = Path(path).expanduser().resolve()
        p.parent.mkdir(parents=True, exist_ok=True)
        
        wb = Workbook()
        wb.remove(wb.active)
        
        for sheet_data in data.get("sheets", []):
            ws = wb.create_sheet(title=sheet_data.get("name", "Hoja1"))
            
            header_color = sheet_data.get("header_color", "4472C4")
            header_fill = PatternFill(start_color=header_color, end_color=header_color, fill_type="solid")
            header_font = Font(bold=True, color="FFFFFF")
            thin_border = Border(
                left=Side(style="thin"), right=Side(style="thin"),
                top=Side(style="thin"), bottom=Side(style="thin")
            )
            
            row_offset = sheet_data.get("start_row", 1)
            
            # Encabezados
            headers = sheet_data.get("headers", [])
            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=row_offset, column=col_idx, value=header)
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal="center")
                cell.border = thin_border
            
            # Filas de datos (soporta fórmulas con "=")
            for r, row in enumerate(sheet_data.get("rows", []), row_offset + 1):
                for c, val in enumerate(row, 1):
                    cell = ws.cell(row=r, column=c)
                    if isinstance(val, str) and val.startswith("="):
                        cell.value = val  # openpyxl guarda fórmulas nativamente
                    elif isinstance(val, (int, float)):
                        cell.value = val
                    else:
                        cell.value = str(val) if val is not None else ""
                    cell.border = thin_border
            
            # Fórmulas adicionales en celdas específicas
            for cell_ref, formula in sheet_data.get("formulas", {}).items():
                ws[cell_ref] = formula
                ws[cell_ref].font = Font(bold=True)
            
            # Anchos de columna
            for col, width in sheet_data.get("column_widths", {}).items():
                ws.column_dimensions[col].width = width
            
            # Auto-ajustar si no se especificaron anchos
            if not sheet_data.get("column_widths"):
                for col_idx in range(1, len(headers) + 1):
                    ws.column_dimensions[get_column_letter(col_idx)].width = 15
        
        wb.save(str(p))
        return json.dumps({
            "success": True, "path": str(p), "type": "xlsx",
            "sheets": len(data.get("sheets", []))
        }, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": f"Error creando Excel: {e}"})


@mcp.tool()
def pptx_create(path: str, content: str) -> str:
    """
    Crea una presentación PowerPoint (.pptx) completa.
    
    Args:
        path: Ruta del archivo .pptx a crear
        content: JSON string con la estructura:
            {
              "title": "Título de la presentación",
              "subtitle": "Subtítulo / autor",
              "slides": [
                {"title": "Diapositiva 1",
                 "bullets": ["Punto 1", "Punto 2", "Punto 3"],
                 "notes": "Notas del presentador (opcional)"},
                {"title": "Con tabla",
                 "table": {"headers": ["A","B"], "rows": [["1","2"],["3","4"]]}},
                {"title": "Solo texto", "text": "Párrafo de texto libre"}
              ]
            }
    """
    if not PPTX_AVAILABLE:
        return json.dumps({"error": "pip install python-pptx"})
    
    try:
        data = json.loads(content) if isinstance(content, str) else content
        p = Path(path).expanduser().resolve()
        p.parent.mkdir(parents=True, exist_ok=True)
        
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)  # 16:9
        prs.slide_height = PptxInches(7.5)
        
        # Diapositiva de título
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        if slide.shapes.title:
            slide.shapes.title.text = data.get("title", "Presentación")
        if len(slide.placeholders) > 1 and data.get("subtitle"):
            slide.placeholders[1].text = data["subtitle"]
        
        # Diapositivas de contenido
        content_layout = prs.slide_layouts[1]  # Título + contenido
        blank_layout = prs.slide_layouts[6]    # En blanco
        
        for slide_data in data.get("slides", []):
            has_table = bool(slide_data.get("table"))
            slide = prs.slides.add_slide(blank_layout if has_table else content_layout)
            
            # Título
            if has_table:
                # Agregar cuadro de título manual en layout en blanco
                from pptx.util import Emu
                txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3),
                                                  PptxInches(12), PptxInches(1))
                tf = txBox.text_frame
                tf.text = slide_data.get("title", "")
                tf.paragraphs[0].font.size = PptxPt(32)
                tf.paragraphs[0].font.bold = True
            else:
                if slide.shapes.title:
                    slide.shapes.title.text = slide_data.get("title", "")
            
            # Bullets
            if slide_data.get("bullets") and not has_table:
                body = slide.placeholders[1].text_frame
                body.clear()
                for i, bullet in enumerate(slide_data["bullets"]):
                    para = body.paragraphs[0] if i == 0 else body.add_paragraph()
                    if isinstance(bullet, dict):
                        para.text = str(bullet.get("text", ""))
                        para.level = bullet.get("level", 0)
                    else:
                        para.text = str(bullet)
            
            # Texto libre
            if slide_data.get("text") and not has_table:
                body = slide.placeholders[1].text_frame
                body.text = slide_data["text"]
            
            # Tabla
            if has_table:
                tdata = slide_data["table"]
                headers = tdata.get("headers", [])
                rows = tdata.get("rows", [])
                if headers:
                    shape = slide.shapes.add_table(
                        len(rows) + 1, len(headers),
                        PptxInches(0.5), PptxInches(1.5),
                        PptxInches(12), PptxInches(0.8 * (len(rows) + 1))
                    )
                    table = shape.table
                    for c, h in enumerate(headers):
                        table.cell(0, c).text = str(h)
                    for r, row in enumerate(rows):
                        for c, val in enumerate(row):
                            if c < len(headers):
                                table.cell(r + 1, c).text = str(val)
            
            # Notas del presentador
            if slide_data.get("notes"):
                slide.notes_slide.notes_text_frame.text = slide_data["notes"]
        
        prs.save(str(p))
        return json.dumps({
            "success": True, "path": str(p), "type": "pptx",
            "slides": len(data.get("slides", [])) + 1
        }, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": f"Error creando PowerPoint: {e}"})


@mcp.tool()
def pdf_create(path: str, content: str) -> str:
    """
    Crea un PDF profesional, con soporte para MEMBRETE (logo, encabezado, pie).
    
    Args:
        path: Ruta del archivo .pdf a crear
        content: JSON string con la estructura:
            {
              "letterhead": {
                  "company": "Mi Empresa S.A.",
                  "tagline": "Eslogan (opcional)",
                  "logo_path": "C:/ruta/logo.png (opcional)",
                  "color": "#1a5276",
                  "footer": "Texto del pie de página (opcional)"
              },
              "title": "Título del documento",
              "date": "2026-07-28 (opcional)",
              "sections": [
                {"heading": "1. Introducción", "text": "Contenido..."},
                {"heading": "2. Datos", "bullets": ["a", "b"]},
                {"table": {"headers": ["Col1","Col2"], "rows": [["x","y"]]}}
              ]
            }
    """
    if not FPDF_AVAILABLE:
        return json.dumps({"error": "pip install fpdf2"})
    
    try:
        data = json.loads(content) if isinstance(content, str) else content
        p = Path(path).expanduser().resolve()
        p.parent.mkdir(parents=True, exist_ok=True)
        
        letterhead = data.get("letterhead", {})
        color_hex = letterhead.get("color", "#1a5276").lstrip("#")
        r, g, b = tuple(int(color_hex[i:i+2], 16) for i in (0, 2, 4))
        
        class LetterheadPDF(FPDF):
            def header(self):
                if letterhead.get("logo_path") and Path(letterhead["logo_path"]).exists():
                    self.image(letterhead["logo_path"], 10, 8, 25)
                self.set_font("helvetica", "B", 16)
                self.set_text_color(r, g, b)
                if letterhead.get("company"):
                    self.cell(0, 8, letterhead["company"], align="C", new_x="LMARGIN", new_y="NEXT")
                if letterhead.get("tagline"):
                    self.set_font("helvetica", "I", 10)
                    self.set_text_color(100, 100, 100)
                    self.cell(0, 6, letterhead["tagline"], align="C", new_x="LMARGIN", new_y="NEXT")
                # Línea separadora del membrete
                self.set_draw_color(r, g, b)
                self.set_line_width(0.8)
                self.line(10, self.get_y() + 2, 200, self.get_y() + 2)
                self.ln(8)
            
            def footer(self):
                self.set_y(-15)
                self.set_font("helvetica", "I", 8)
                self.set_text_color(128, 128, 128)
                footer_text = letterhead.get("footer", "")
                if footer_text:
                    self.cell(0, 5, footer_text, align="C", new_x="LMARGIN", new_y="NEXT")
                self.cell(0, 5, f"Página {self.page_no()}/{{nb}}", align="C")
        
        pdf = LetterheadPDF()
        pdf.set_auto_page_break(auto=True, margin=20)
        pdf.add_page()
        
        # Título del documento
        if data.get("title"):
            pdf.set_font("helvetica", "B", 20)
            pdf.set_text_color(0, 0, 0)
            pdf.multi_cell(0, 10, data["title"], new_x="LMARGIN", new_y="NEXT")
            pdf.ln(2)
        
        if data.get("date"):
            pdf.set_font("helvetica", "", 10)
            pdf.set_text_color(100, 100, 100)
            pdf.cell(0, 6, f"Fecha: {data['date']}", new_x="LMARGIN", new_y="NEXT")
            pdf.ln(4)
        
        # Secciones
        for section in data.get("sections", []):
            if section.get("heading"):
                pdf.set_font("helvetica", "B", 14)
                pdf.set_text_color(r, g, b)
                pdf.multi_cell(0, 8, section["heading"], new_x="LMARGIN", new_y="NEXT")
                pdf.ln(1)
            
            if section.get("text"):
                pdf.set_font("helvetica", "", 11)
                pdf.set_text_color(30, 30, 30)
                pdf.multi_cell(0, 6, section["text"], new_x="LMARGIN", new_y="NEXT")
                pdf.ln(3)
            
            for bullet in section.get("bullets", []):
                pdf.set_font("helvetica", "", 11)
                pdf.set_text_color(30, 30, 30)
                # new_x=LMARGIN es OBLIGATORIO: sin él, fpdf2 deja el cursor
                # al margen derecho y el siguiente multi_cell falla con
                # "Not enough horizontal space to render a single character"
                pdf.multi_cell(0, 6, f"  - {bullet}", new_x="LMARGIN", new_y="NEXT")
            if section.get("bullets"):
                pdf.ln(3)
            
            table_data = section.get("table")
            if table_data and table_data.get("headers"):
                headers = table_data["headers"]
                rows = table_data.get("rows", [])
                col_width = 190 / len(headers)
                
                pdf.set_font("helvetica", "B", 10)
                pdf.set_fill_color(r, g, b)
                pdf.set_text_color(255, 255, 255)
                for h in headers:
                    pdf.cell(col_width, 8, str(h), border=1, fill=True, align="C")
                pdf.ln()
                
                pdf.set_font("helvetica", "", 10)
                pdf.set_text_color(30, 30, 30)
                for row in rows:
                    for c, val in enumerate(row):
                        if c < len(headers):
                            pdf.cell(col_width, 7, str(val), border=1, align="C")
                    pdf.ln()
                pdf.ln(4)
        
        pdf.output(str(p))
        return json.dumps({"success": True, "path": str(p), "type": "pdf"}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": f"Error creando PDF: {e}"})


# ═════════════════════════════════════════════════════════════════════════════
# 10. LECTURA WEB (fetch integrado)
# ═════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def web_fetch(url: str, max_chars: int = 8000) -> str:
    """
    Lee una página web y devuelve su contenido como texto/Markdown limpio.
    Útil para que la IA lea documentación, noticias o artículos.
    
    Args:
        url: URL de la página (http:// o https://)
        max_chars: Máximo de caracteres a devolver (default: 8000)
    """
    if not REQUESTS_AVAILABLE:
        return json.dumps({"error": "pip install requests"})
    
    try:
        import re
        from html import unescape
        
        resp = _requests.get(
            url,
            timeout=20,
            headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) MCP-Windows-AI"},
        )
        resp.raise_for_status()
        html = resp.text
        
        # Convertir a Markdown si markdownify está disponible
        if MARKDOWNIFY_AVAILABLE:
            text = _md(html, heading_style="ATX", strip=["script", "style"])
        else:
            # Fallback: limpieza básica con regex
            text = re.sub(r"(?is)<(script|style|noscript)[^>]*>.*?</\1>", " ", html)
            text = re.sub(r"(?s)<[^>]+>", " ", text)
            text = unescape(text)
            text = re.sub(r"\s+", " ", text).strip()
        
        truncated = len(text) > max_chars
        return json.dumps({
            "url": url,
            "status": resp.status_code,
            "length": len(text),
            "truncated": truncated,
            "content": text[:max_chars],
        }, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": f"No se pudo leer {url}: {e}"})


# ═════════════════════════════════════════════════════════════════════════════
# 11. APERTURA E INTERACCIÓN CON PROGRAMAS
# ═════════════════════════════════════════════════════════════════════════════

# Programas conocidos con sus rutas/comandos comunes
KNOWN_PROGRAMS = {
    "notepad": "notepad.exe",
    "bloc de notas": "notepad.exe",
    "paint": "mspaint.exe",
    "calculadora": "calc.exe",
    "explorer": "explorer.exe",
    "explorador": "explorer.exe",
    "cmd": "cmd.exe",
    "powershell": "powershell.exe",
    "word": "winword.exe",
    "excel": "excel.exe",
    "powerpoint": "powerpnt.exe",
    "chrome": "chrome.exe",
    "edge": "msedge.exe",
    "firefox": "firefox.exe",
    "vscode": "code",
    "visual studio code": "code",
    "spotify": "spotify.exe",
    "vlc": "vlc.exe",
    "coreldraw": "CorelDRW.exe",
    "photoshop": "Photoshop.exe",
    "illustrator": "Illustrator.exe",
    "gimp": "gimp.exe",
    "inkscape": "inkscape.exe",
    "blender": "blender.exe",
    "obs": "obs64.exe",
    "git bash": "git-bash.exe",
    "terminal": "wt.exe",
    "windows terminal": "wt.exe",
}

@mcp.tool()
def open_program(program: str, file_path: str = "", arguments: str = "") -> str:
    """
    Abre un programa, opcionalmente con un archivo o argumentos.
    Reconoce nombres comunes: word, excel, powerpoint, chrome, vscode, notepad,
    paint, calculadora, coreldraw, photoshop, gimp, inkscape, blender, vlc, etc.
    
    Args:
        program: Nombre del programa (ej: "excel", "vscode", "chrome")
        file_path: Archivo a abrir con el programa (opcional)
        arguments: Argumentos extra de línea de comandos (opcional)
    """
    try:
        prog_lower = program.lower().strip()
        executable = KNOWN_PROGRAMS.get(prog_lower, program)
        
        cmd_parts = [executable]
        if file_path:
            p = Path(file_path).expanduser().resolve()
            if not p.exists():
                return json.dumps({"error": f"El archivo no existe: {file_path}"})
            cmd_parts.append(str(p))
        if arguments:
            cmd_parts.extend(arguments.split())
        
        # Intentar con 'start' de Windows para programas instalados en rutas del sistema
        try:
            subprocess.Popen(cmd_parts, shell=True)
        except FileNotFoundError:
            # Fallback: usar 'start' que busca en el registro de Windows
            start_cmd = f'start "" "{executable}"'
            if file_path:
                start_cmd += f' "{file_path}"'
            if arguments:
                start_cmd += f' {arguments}'
            subprocess.Popen(start_cmd, shell=True)
        
        return json.dumps({
            "success": True,
            "program": program,
            "executable": executable,
            "file": file_path or None,
        }, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": f"No se pudo abrir '{program}': {e}"})


@mcp.tool()
def list_installed_programs() -> str:
    """
    Lista los programas instalados en Windows (desde el registro).
    Útil para saber qué software de diseño/desarrollo está disponible.
    """
    if not WINREG_AVAILABLE:
        return json.dumps({"error": "winreg no disponible"})
    
    programs = []
    registry_paths = [
        (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\Microsoft\Windows\CurrentVersion\Uninstall"),
        (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\WOW6432Node\Microsoft\Windows\CurrentVersion\Uninstall"),
        (winreg.HKEY_CURRENT_USER, r"SOFTWARE\Microsoft\Windows\CurrentVersion\Uninstall"),
    ]
    
    for hive, path in registry_paths:
        try:
            with winreg.OpenKey(hive, path) as key:
                i = 0
                while True:
                    try:
                        subkey_name = winreg.EnumKey(key, i)
                        with winreg.OpenKey(key, subkey_name) as subkey:
                            try:
                                name, _ = winreg.QueryValueEx(subkey, "DisplayName")
                                try:
                                    version, _ = winreg.QueryValueEx(subkey, "DisplayVersion")
                                except FileNotFoundError:
                                    version = ""
                                if name and name not in [p["name"] for p in programs]:
                                    programs.append({"name": name, "version": version})
                            except FileNotFoundError:
                                pass
                        i += 1
                    except OSError:
                        break
        except FileNotFoundError:
            continue
    
    programs.sort(key=lambda x: x["name"].lower())
    return json.dumps({
        "count": len(programs),
        "programs": programs,
    }, ensure_ascii=False, indent=2)


@mcp.tool()
def type_in_program(program_title: str, text: str, delay: float = 1.0) -> str:
    """
    Enfoca una ventana de un programa y escribe texto en ella.
    Combina focus_window + keyboard_type en un solo paso.
    
    Args:
        program_title: Título (o parte) de la ventana del programa
        text: Texto a escribir
        delay: Segundos de espera tras enfocar antes de escribir
    """
    if not PYAUTOGUI_AVAILABLE:
        return json.dumps({"error": "Se requiere pyautogui"})
    
    try:
        # Enfocar ventana
        if PYWINDOW_AVAILABLE:
            windows = gw.getWindowsWithTitle(program_title)
            if not windows:
                return json.dumps({"error": f"No se encontró ventana: {program_title}"})
            win = windows[0]
            if win.isMinimized:
                win.restore()
            win.activate()
            time.sleep(delay)
        
        # Escribir
        pyautogui.typewrite(text, interval=0.02)
        return json.dumps({
            "success": True,
            "window": program_title,
            "chars_written": len(text),
        }, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e)})


# ═════════════════════════════════════════════════════════════════════════════
# C10 · CAPA DE EJECUCION PROFESIONAL (SPEC C1)
# screen_capture · read_ui_state · ocr_screen · open_app
# ═════════════════════════════════════════════════════════════════════════════

def _c10_capture_region(left=0, top=0, width=None, height=None):
    """Captura la pantalla o una region. Returns (PIL.Image, path)."""
    from PIL import ImageGrab
    if width and height:
        img = ImageGrab.grab(bbox=(left, top, left + width, top + height))
    else:
        img = ImageGrab.grab()
    return img


@mcp.tool()
def screen_capture(region: str = "full", filename: str = "") -> str:
    """
    Captura la pantalla (o una region) como PNG. C10 de ejecucion verificada.

    Args:
        region: "full" (toda la pantalla) o "X,Y,W,H" (left,top,width,height)
        filename: nombre opcional para el archivo (default: screen_YYYYMMDD_HHMMSS.png)

    Returns:
        Ruta del archivo PNG
    """
    try:
        from PIL import ImageGrab
        if region == "full":
            img = ImageGrab.grab()
        else:
            try:
                parts = [int(x) for x in region.split(",")]
                if len(parts) != 4:
                    return json.dumps({"error": "region debe ser 'full' o 'left,top,width,height'"})
                l, t, w, h = parts
                img = ImageGrab.grab(bbox=(l, t, l + w, t + h))
            except Exception as e:
                return json.dumps({"error": f"region invalida: {e}"})

        out_dir = Path.home() / "atlas_shots"
        out_dir.mkdir(exist_ok=True)
        fname = filename or f"screen_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png"
        filepath = str(out_dir / fname)
        img.save(filepath)
        return json.dumps({"success": True, "path": filepath, "size": img.size},
                          ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def read_ui_state(window_title: str = "", max_depth: int = 4) -> str:
    """
    Lee el arbol UIA (UI Automation) de la ventana activa o la indicada,
    devolviendo JSON con los controles visibles (type, name, automation_id, bounds).

    Args:
        window_title: titulo (o parte) de la ventana; vacio = ventana activa
        max_depth: profundidad maxima del arbol (default 4)

    Returns:
        JSON con el arbol de controles
    """
    try:
        import uiautomation as uia
    except Exception as e:
        return json.dumps({"error": f"UIA no disponible: {e}"})

    try:
        win = None
        if window_title:
            win = uia.WindowControl(searchDepth=1, SubName=window_title)
            if not win.Exists(0.5, 0.1):
                return json.dumps({"error": f"ventana '{window_title}' no encontrada",
                                   "window": window_title})
        else:
            win = uia.GetForegroundControl()

        def walk(control, depth):
            if depth > max_depth:
                return None
            node = {}
            try:
                node["type"] = control.ControlTypeName
            except Exception:
                node["type"] = "?"
            try:
                node["name"] = control.Name
            except Exception:
                node["name"] = ""
            try:
                node["automation_id"] = control.AutomationId
            except Exception:
                node["automation_id"] = ""
            try:
                r = control.BoundingRectangle
                node["bounds"] = [int(r.left), int(r.top),
                                  int(r.right - r.left), int(r.bottom - r.top)]
            except Exception:
                node["bounds"] = None
            children = []
            try:
                for c in control.GetChildren():
                    ch = walk(c, depth + 1)
                    if ch:
                        children.append(ch)
            except Exception:
                pass
            if children:
                node["children"] = children[:20]
            return node

        root = walk(win, 0)
        return json.dumps({"success": True, "window": window_title or "active", "tree": root},
                          ensure_ascii=False, default=str)
    except Exception as e:
        return json.dumps({"error": f"fallo UIA: {e}"})


@mcp.tool()
def ocr_screen(region: str = "full", language: str = "es") -> str:
    """
    Extrae texto de la pantalla (o region) usando OCR nativo de Windows (WinRT).

    Args:
        region: "full" o "left,top,width,height"
        language: idioma del OCR (default "es"; usar "en" si falla)

    Returns:
        JSON con el texto extraido y su distribucion por bloques
    """
    try:
        from PIL import ImageGrab
        import winsdk.windows.media.ocr as ocr_ws
        import winsdk.windows.globalization as glob_ws
        import winsdk.windows.graphics.imaging as imaging_ws
        import winsdk.windows.storage.streams as streams_ws
    except Exception as e:
        return json.dumps({"error": f"WinRT OCR no disponible: {e}"})

    try:
        if region == "full":
            img = ImageGrab.grab()
        else:
            parts = [int(x) for x in region.split(",")]
            l, t, w, h = parts
            img = ImageGrab.grab(bbox=(l, t, l + w, t + h))

        import io
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        data = buf.getvalue()

        import winsdk.windows.storage.streams as streams
        import winsdk.windows.storage as storage
        import winsdk.windows.foundation as foundation

        # Escribir temporal para WinRT
        tmp = Path.home() / "atlas_shots"
        tmp.mkdir(exist_ok=True)
        tmp_file = tmp / "ocr_input.png"
        tmp_file.write_bytes(data)

        import asyncio

        async def _run_ocr():
            file = await storage.StorageFile.get_file_from_path_async(str(tmp_file))
            stream = await file.open_async(storage.FileAccessMode.READ)
            decoder = await imaging_ws.BitmapDecoder.create_async(stream)
            bitmap = await decoder.get_software_bitmap_async()
            engine = ocr_ws.OcrEngine.try_create_from_language(
                glob_ws.Language(language))
            if engine is None:
                engine = ocr_ws.OcrEngine.try_create_from_user_profile_languages()
            result = await engine.recognize_async(bitmap)
            lines = []
            for line in result.lines:
                lines.append({
                    "text": line.text,
                    "words": [w.text for w in line.words],
                })
            return lines

        try:
            lines = asyncio.run(_run_ocr())
        except RuntimeError:
            # loop ya existente
            loop = asyncio.get_event_loop()
            future = asyncio.ensure_future(_run_ocr())
            lines = loop.run_until_complete(future)

        return json.dumps({"success": True, "lines": lines,
                           "text": "\n".join(l["text"] for l in lines)},
                          ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": f"fallo OCR: {e}"})


@mcp.tool()
def open_app(program: str, args: str = "", path: str = "") -> str:
    """
    Abre un programa registrado con la validacion del guardian.

    Args:
        program: nombre o ruta del programa (ej: "chrome", "vscode", "notepad")
        args: argumentos adicionales
        path: ruta exacta del ejecutable (opcional; usa registro si no)

    Returns:
        Confirmacion con PID del proceso
    """
    block = _guardian_gate("process_start", {"command": program, "args": args})
    if block:
        return block
    try:
        import subprocess
        if path:
            cmd = f'"{path}" {args}'.strip()
        else:
            # buscar en registro si open_program lo conoce
            known = {
                "chrome": r"C:\Program Files\Google\Chrome\Application\chrome.exe",
                "edge": r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
                "notepad": "notepad.exe",
                "vscode": r"C:\Users\Administrator\AppData\Local\Programs\Microsoft VS Code\Code.exe",
                "cmd": "cmd.exe",
                "powershell": "powershell.exe",
                "explorer": "explorer.exe",
                "paint": "mspaint.exe",
                "calc": "calc.exe",
                "word": r"C:\Program Files\Microsoft Office\root\Office16\WINWORD.EXE",
                "excel": r"C:\Program Files\Microsoft Office\root\Office16\EXCEL.EXE",
                "corel": r"C:\Program Files\Corel\CorelDRAW Graphics Suite\Programs\CorelDRAW.exe",
            }
            exe = known.get(program.lower(), program)
            cmd = f'"{exe}" {args}'.strip() if (" " in exe) else f"{exe} {args}".strip()
        proc = subprocess.Popen(cmd, shell=False)
        return json.dumps({"success": True, "program": program, "pid": proc.pid},
                          ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e)})


# ═════════════════════════════════════════════════════════════════════════════
# PUNTO DE ENTRADA
# ═════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    print("=" * 60)
    print("  MCP Windows Automation Server")
    print("  Ejecutando... Esperando conexion del cliente MCP")
    print("=" * 60)
    print("")
    print("  Tools cargadas:")
    print("  [W] Ventanas:     list, focus, move, resize, min/max, close")
    print("  [M] Mouse:        move, click, scroll, position, drag")
    print("  [K] Teclado:      type, hotkey, press")
    print("  [F] Archivos:     list, read, write, delete, copy, folder_create")
    print("  [P] Procesos:     list, kill, start")
    print("  [I] Sistema:      info, screenshot, clipboard, open_url")
    print("  [C] Comandos:     run_command, volume_set, wifi_info")
    print("  [R] Registry:     registry_read")
    print("  [N] Notify:       notify")
    print("  [D] Documentos:   docx_create, xlsx_create, pptx_create, pdf_create")
    print("  [A] Apps:         open_program, list_installed_programs, type_in_program")
    print("  [S] Scripts:      run_script")
    print("  [C10] C1:         screen_capture, read_ui_state, ocr_screen, open_app")
    print("")
    
    mcp.run(transport="stdio")
